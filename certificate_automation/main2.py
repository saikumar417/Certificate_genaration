# main2.py
# AUTHOR: N SAI KUMAR, Registration.No: 23751F0020
# Guidance: Mr. J. SHEIK MOHAMED Sir CO-ORDINATER,SITAMS.
# BRANCH: MCA 2023-25 Batch
# Taken help from OpenAI's GPT-4 architecture. It helped me more.

import pandas as pd
from PIL import Image, ImageDraw, ImageFont
import os
import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.utils import ImageReader
from pptx.util import Inches
from watermark_module import *
from permissions import *
import copy

import sys
import time
import email_module

seminar_title = ""


# Variables to hold values set by the UI
knowledge_text = ""
passion_text = ""
template = ""
logo1 = ""
logo2 = ""
logo3 = ""
logo4 = ""
  # This is assigned in the UI like: main2.seminar_title = ...

def generate_certificates():
    global knowledge_text, passion_text, template, logo1, logo2, logo3, logo4,seminar_title  # Get text values from UI
    global x1, x2, x3, x4, y1, y2, y3, y4, size1, size2, size3, size4  # Get x and y values and size values

    print("Knowledge Text:", knowledge_text)
    print("Passion Text:", passion_text)
    print("Template Path:", template)
    print("seminar_title:", seminar_title)
    print("Logo 1 Path:", logo1, "at (", x1, ",", y1, ") and size:", size1)
    print("Logo 2 Path:", logo2, "at (", x2, ",", y2, ") and size:", size2)
    print("Logo 3 Path:", logo3, "at (", x3, ",", y3, ") and size:", size3)
    print("Logo 4 Path:", logo4, "at (", x4, ",", y4, ") and size:", size4)

    # Certificate generation logic here...

    # Ensure values are received correctly
    if not template:
        print("Error: Template is missing")
        return

    # Add more debug prints where x and y are calculated or used
    x = 100  # Example
    y = 200  # Example
    print(f"x: {x}, y: {y}")

    # Add certificate generation logic here

    # Function to add text to image
    def add_text_to_image(image, text, position, font, color):
        draw = ImageDraw.Draw(image)
        draw.text(position, text, font=font, fill=color)

    # Function to select the Excel file
    def select_excel_file():
        root = tk.Tk()
        root.withdraw()  # Hide the root window
        file_path = filedialog.askopenfilename(title="Select Excel File", filetypes=[("Excel files", "*.xlsx")])
        return file_path

    # Ask the user to select the Excel file
    excel_file_path = select_excel_file()

    # Ensure a file was selected
    if not excel_file_path:
        print("No file selected. Exiting...")
        exit()

    # Automatically set the output directory based on the Excel file's directory
    output_dir = os.path.join(os.path.dirname(excel_file_path), 'certificates_output')
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    # Load Excel data
    df = pd.read_excel(excel_file_path)

    # Get the directory where the current Python script is located
    current_dir = os.path.dirname(os.path.abspath(__file__))

    # Define the folder where the input files (template and logos) are stored
    input_folder = os.path.join(current_dir, 'input_files')

    # Load the certificate template and logos from the 'input_files' folder
    template_image = Image.open(template)
    logo1_img = None
    logo2_img = None
    logo3_img = None
    logo4_img = None
    if logo1:
        logo1_img = Image.open(logo1)
    if logo2:
        logo2_img = Image.open(logo2)
    if logo3:
        logo3_img = Image.open(logo3)
    if logo4:
        logo4_img = Image.open(logo4)
    

    # Load the font
    font_path = 'input_files/arial.ttf'  # Path to your TTF file
    font_path2 = 'input_files/Roboto-Italic.ttf'

    #font_path2 = os.path.join(input_folder, 'Roboto-italic.ttf')

    font_large = ImageFont.truetype(font_path, 40)  # For title and header
    font_medium = ImageFont.truetype(font_path, 40)  # For participant info
    font_small = ImageFont.truetype(font_path, 35)  # For additional info
    font_signatures = ImageFont.truetype(font_path, 30)  # For signatures
    font_name_designation = ImageFont.truetype(font_path, 35)  # For names and designations
    font_italic = ImageFont.truetype(font_path2, 35)  # For italic text

    # Get the base name of the input Excel file (without extension)
    excel_base_name = os.path.splitext(os.path.basename(excel_file_path))[0]
    
    # This output path for sending eamil it will excel file with name,email,certificate path
    output_excel_path = os.path.join(output_dir, f"{excel_base_name}_email.xlsx")  # Save output file in the same folder

    certificate_details = []
    # Prepare PDF canvas and PowerPoint presentation using the base name
    pdf = canvas.Canvas(os.path.join(output_dir, f"{excel_base_name}.pdf"), pagesize=A4)
    presentation = Presentation()
    presentation.slide_width = Inches(11.7)
    presentation.slide_height = Inches(8.27)

    # Create a landscape PDF for certificates with the base name
    landscape_pdf_path = os.path.join(output_dir, f"{excel_base_name}.pdf")
    landscape_pdf = canvas.Canvas(landscape_pdf_path, pagesize=landscape(A4))

    #   Title 
    # Function to add wrapped text to image and return the total height used
    def add_wrapped_text_to_image(image, text, position, font, color, max_width):
        draw = ImageDraw.Draw(image)

        # Split the text into lines
        lines = []
        words = text.split()
        current_line = ""

        # Fixed position for the first line
        x_position_first = position[0]  # X position for the first line
        y_position_first = position[1]   # Y position for the first line

        for word in words:
            test_line = current_line + word + " "
            # Calculate the bounding box of the line
            line_bbox = draw.textbbox((0, 0), test_line, font=font)
            line_width = line_bbox[2] - line_bbox[0]  # Calculate the width of the line

            # If the line exceeds the max width, add it to lines and start a new line
            if line_width > max_width:
                lines.append(current_line)
                current_line = word + " "
            else:
                current_line = test_line

        # Append the last line
        if current_line:
            lines.append(current_line)

        # Draw the first line on the image
        if lines:
            draw.text((x_position_first, y_position_first), lines[0], font=font, fill=color)

        # Declare variables for the second line
        x_position_second = x_position_first+20  # Same x position as the first line
        y_position_second = y_position_first+10 + draw.textbbox((0, 0), lines[0], font=font)[3] - draw.textbbox((0, 0), lines[0], font=font)[1]  # Calculate Y position for the second line

        # Draw the second line if it exists
        if len(lines) > 1:
            draw.text((x_position_second, y_position_second), lines[1], font=font, fill=color)

        # Return the height of all drawn lines
        total_height = (draw.textbbox((0, 0), lines[0], font=font)[3] - draw.textbbox((0, 0), lines[0], font=font)[1])
        if len(lines) > 1:
            total_height += (draw.textbbox((0, 0), lines[1], font=font)[3] - draw.textbbox((0, 0), lines[1], font=font)[1])

        return total_height
    #title end
    title = seminar_title if seminar_title else ""

    email_column_exists = 'Email' in df.columns

    print("Request Processing.....")
    # Iterate over each row in the Excel file
    for index, row in df.iterrows():
        crtno = row['Certno'] if pd.notnull(row['Certno']) else ""
    # r_person=row['Resource_person']
        name = row['Name'] if pd.notnull(row['Name']) else ""
        department = row['Department'] if pd.notnull(row['Department']) else ""
        #title = row['Title'] if pd.notnull(row['Title']) else ""
        
        # Ensure 'Date' is a string and format it if needed
        date = row['Date'] if pd.notnull(row["Date"]) else ""
         # Process 'Email' only if the column exists
        if email_column_exists:
            email = row['Email'] if pd.notnull(row['Email']) else ""
        else:
            email = ""
        
        if isinstance(date, pd.Timestamp):  # If it's a datetime object, format it
            date = date.strftime('%d-%m-%Y')
        elif isinstance(date, float):  # If it's a float (probably NaN or missing), set it to an empty string
            date = ""
        elif not isinstance(date, str):  # Convert any other type to a string
            date = str(date)

        ''''if not isinstance(date, str):
            date = date.strftime('%d-%m-%Y')'''

        organized_by = row['Organized_by'] if pd.notnull(row['Organized_by']) else ""

        # Signature paths, names, and designations from UI
        signatures = []
        for i in range(1, 5):
            key = f'signature{i}'
            
            # Only include if the signature path exists and is valid
            path = globals().get(key, "")
            names = globals().get(f'name{i}', "")
            designation = globals().get(f'designation{i}', "")

            if isinstance(path, str) and path.strip() != "" and not pd.isna(path):
                signatures.append({
                    'path': path,
                    'name': names if pd.notnull(names) else "",
                    'designation': designation if pd.notnull(designation) else ""
                })

        # Filter out invalid/empty signatures (preserving your logic)
        valid_signatures = [
            s for s in signatures
            if isinstance(s['path'], str) and s['path'].strip() != "" and not pd.isna(s['path'])
        ]



        # Copy the certificate template to avoid modifying the original
        certificate = template_image.copy()
        #certificate = copy.deepcopy(template)

        if logo1_img:
            logo1_img = logo1_img.convert("RGBA")
        if logo2_img:
            logo2_img = logo2_img.convert("RGBA")
        if logo3_img:
            logo3_img = logo3_img.convert("RGBA")
        if logo4_img:
            logo4_img = logo4_img.convert("RGBA")

        # Resize logos only if they exist
        resized_logo1 = logo1_img.resize((size1, size1), Image.Resampling.LANCZOS) if logo1_img else None
        resized_logo2 = logo2_img.resize((size2, size2), Image.Resampling.LANCZOS) if logo2_img else None
        resized_logo3 = logo3_img.resize((size3, size3), Image.Resampling.LANCZOS) if logo3_img else None
        resized_logo4 = logo4_img.resize((size4, size4), Image.Resampling.LANCZOS) if logo4_img else None

        # Define positions for logos (adjusted based on template size)
        inc_y_pos = 150
        logo1_position = (x1, y1)    # Top-left corner for logo1
        logo2_position = (x2, y2)  # Top-right corner for logo2
        logo3_position = (x3, y3)  # Top-right corner for logo3
        logo4_position = (x4, y4)  # Adjusted to avoid overlap with logo1

        # Paste logos onto the certificate
        if resized_logo1:
            certificate.paste(resized_logo1, logo1_position, resized_logo1)
        if resized_logo2:
            certificate.paste(resized_logo2, logo2_position, resized_logo2)
        if resized_logo3:
            certificate.paste(resized_logo3, logo3_position, resized_logo3)
        if resized_logo4:
            certificate.paste(resized_logo4, logo4_position, resized_logo4)

        
        # Define text positions
        cname_position = ((certificate.width // 2) - 770, inc_y_pos) # collage name position
        autonomous_position = ((certificate.width // 2) - 425, 55 + inc_y_pos)
        chittoor_position = ((certificate.width // 2) - 170, 105 + inc_y_pos)
        saa_position = ((certificate.width // 2) - 360, 155 + inc_y_pos)
        Reg_position = ((certificate.width // 2) - 160, 205 + inc_y_pos)  # registration number
      

        # Calculate text width and center it along X-axis
        text_bbox = ImageDraw.Draw(certificate).textbbox((0, 0), knowledge_text, font=font_italic)
        text_width = text_bbox[2] - text_bbox[0]
        ksharing_position = ((certificate.width - text_width) // 2, 245 + inc_y_pos)  # Center X-axis

        text_bbox = ImageDraw.Draw(certificate).textbbox((0, 0), passion_text, font=font_medium)
        text_width = text_bbox[2] - text_bbox[0]
        passion_position = ((certificate.width - text_width) // 2, 295 + inc_y_pos)  # Center X-axis


        cop_position = ((certificate.width // 2) - 240, 380 + inc_y_pos)  # Certificate of Participation

        inc_y_pos = 110
        certno_lable_position = (100, 500 + inc_y_pos)
        certno_value_position = (300, 500 + inc_y_pos)

        label_x_position = 300
        value_x_position = 600
        name_position = (label_x_position, 620 + inc_y_pos)
        department_position = (label_x_position, 680 + inc_y_pos)
        participated_position = (label_x_position, 740 + inc_y_pos)
        date_position = (label_x_position, 800 + inc_y_pos)
        organized_by_position = (label_x_position, 860 + inc_y_pos)
        
        # Add diagonal watermark
        add_name_watermark(certificate, name, font_medium, DEFAULT_WATERMARK_COLOR, DEFAULT_TRANSPARENCY)

        # Add text to the certificate with different font sizes
        add_text_to_image(certificate, "SREENIVASA INSTITUTE OF TECHNOLOGY AND MANAGEMENT STUDIES", cname_position, font_large, "#C91902")
        add_text_to_image(certificate, "(Autonomous) (NBA Accredited - CSE,ECE,EEE,MCA)", autonomous_position, font_small, "#3282F6")
        add_text_to_image(certificate, "CHITTOOR-517127", chittoor_position, font_small, "black")
        add_text_to_image(certificate, "SITAMS ALUMNI ASSOCIATION(SAA)", saa_position, font_medium, "#630857")
        add_text_to_image(certificate, "(Reg.No.)", Reg_position, font_small, "#611703")

        add_text_to_image(certificate, "Cert. No. : ", certno_lable_position, font_medium, "red")
        add_text_to_image(certificate, f"{crtno}", certno_value_position, font_medium, "darkblue")

        #add_text_to_image(certificate, "An Alumni Knowledge Sharing Series", ksharing_position, font_italic, "#023842")
        #add_text_to_image(certificate, '"Passionate Towards Passion - An Approach To Career Path"', passion_position, font_medium, "#4F014F")
        # Only add text if it's not empty and not the placeholder text

        add_text_to_image(certificate, knowledge_text, ksharing_position, font_italic, "#023842")
        add_text_to_image(certificate, passion_text, passion_position, font_medium, "#4F014F")

        add_text_to_image(certificate, "Certificate of Participation", cop_position, font_large, "blue")
        add_text_to_image(certificate, "_____________________", ((certificate.width // 2) - 240, 535), font_large, "blue")

    # add_text_to_image(certificate, f"Resource Person: {r_person}",(300,670),font_small,"#1C4463")
    
        # Participant details
        add_text_to_image(certificate, 'Name', name_position, font_medium, "darkred")
        add_text_to_image(certificate, f': {name}', (value_x_position, name_position[1]), font_medium, "black")
        
        add_text_to_image(certificate, 'Department', department_position, font_medium, "darkred")
        add_text_to_image(certificate, f': {department}', (value_x_position, department_position[1]), font_medium, "black")
        
        # Define maximum width for text wrapping
        max_text_width = certificate.width - value_x_position - 50  # Adjust this width as needed
        
        # Add title with text wrapping
        add_text_to_image(certificate, 'Title', participated_position, font_medium, "darkred")
        title_text = f': {title}'
        # After adding the wrapped title
        title_height = add_wrapped_text_to_image(certificate, title_text, (value_x_position, participated_position[1]), font_medium, "black", max_text_width)

        # Adjust the positions of the "Date" and "Organized by" sections dynamically based on title height
        new_date_position = (label_x_position, participated_position[1] + title_height + 20) 

        new_organized_by_position = (label_x_position, new_date_position[1] + 60)
            
        # Adjust positions dynamically for the title, date, and organized by sections
        title_height = add_wrapped_text_to_image(certificate, title_text, (value_x_position, participated_position[1]), font_medium, "black", max_text_width)

        # Update positions based on the height of the title
        new_date_position = (label_x_position, participated_position[1] + title_height + 30)  # Adding padding after the title
        new_organized_by_position = (label_x_position, new_date_position[1] + 70)  # Adding space after the date

        # Adding Date and Organized by fields
        add_text_to_image(certificate, 'Date', new_date_position, font_medium, "darkred")
        add_text_to_image(certificate, f': {date}', (value_x_position, new_date_position[1]), font_medium, "black")

        add_text_to_image(certificate, 'Organized by', new_organized_by_position, font_medium, "darkred")
        add_text_to_image(certificate, f': {organized_by}', (value_x_position, new_organized_by_position[1]), font_medium, "black")

        # Calculate signature positions based on the number of signatures
        num_signatures = len(valid_signatures)
        
        y_position = certificate.height - 330
        if num_signatures == 2:
            signature_positions = [(150, y_position), (certificate.width - 350, y_position)]
        elif num_signatures == 3:
            signature_positions = [(150, y_position), (certificate.width // 2 - 125, y_position), (certificate.width - 500, y_position)]
        elif num_signatures == 4:
            signature_positions = [(150, y_position), (certificate.width // 2 - 300, y_position), (certificate.width // 2 + 150, y_position), (certificate.width - 300, y_position)]

        # Add signature images, names, and designations
        for i, signature in enumerate(valid_signatures):
            if isinstance(signature['path'], str) and os.path.exists(signature['path']):
                # Load and resize the signature
                sig_image = Image.open(signature['path']).resize((200, 100), Image.Resampling.LANCZOS)
                # Ensure signature images are in RGBA mode to handle transparency
                sig_image = sig_image.convert("RGBA")

                # Paste the signature image with proper transparency handling
                certificate.paste(sig_image, signature_positions[i], sig_image)

                # Calculate the width of the name and designation for centering
                name_text = str(signature['name'])
                designation_text = str(signature['designation'])

                # Use ImageDraw.textbbox to calculate the bounding box of the text
                draw = ImageDraw.Draw(certificate)
                name_bbox = draw.textbbox((0, 0), name_text, font=font_name_designation)
                designation_bbox = draw.textbbox((0, 0), designation_text, font=font_name_designation)

                # Calculate the width of the text
                name_width = name_bbox[2] - name_bbox[0]
                designation_width = designation_bbox[2] - designation_bbox[0]

                # Calculate positions to align the name and designation at the center of the signature
                sig_center_x = signature_positions[i][0] + sig_image.width // 2
                name_position = (sig_center_x - name_width // 2, signature_positions[i][1] + 110)
                designation_position = (sig_center_x - designation_width // 2, signature_positions[i][1] + 150)

                # Add the name and designation centered under the signature
                add_text_to_image(certificate, name_text, name_position, font_name_designation, "#0B4002")
                add_text_to_image(certificate, designation_text, designation_position, font_name_designation, "black")

        # Save each certificate as an individual PNG file
        output_path = os.path.join(output_dir, f'{index + 1}_{name}.png')
        certificate.save(output_path)

        if email_column_exists:    # Store the details in the list
            certificate_details.append((name, email, output_path))

        # Add each saved certificate image to the landscape PDF
        landscape_pdf.drawImage(ImageReader(output_path), 0, 0, width=landscape(A4)[0], height=landscape(A4)[1])
        landscape_pdf.showPage()  # Finish the current page in the PDF

        # Add certificate to PowerPoint
        slide_layout = presentation.slide_layouts[6]  # Blank slide layout
        slide = presentation.slides.add_slide(slide_layout)
        slide.shapes.add_picture(output_path, 0, 0, width=presentation.slide_width, height=presentation.slide_height)


    # Save all email details at once after iteration
    if email_column_exists:
        email_module.save_email_details(output_excel_path, certificate_details)
    # Save the landscape PDF
    landscape_pdf.save()

    # Secure Certificate
    input_pdf = landscape_pdf_path
    output_pdf = os.path.join(output_dir, f'{excel_base_name}.pdf')
    protect_pdf(input_pdf, output_pdf)

    # Save the PowerPoint presentation with the base name
    output_ppt_path = os.path.join(output_dir, f"{excel_base_name}.pptx")
    presentation.save(output_ppt_path)

    print(f"All certificates saved in a single PDF: {os.path.join(output_dir, f'{excel_base_name}.pdf')}")
    print(f"All certificates saved in a single landscape PDF: {landscape_pdf_path}")
    print(f"All certificates saved in a single PowerPoint: {output_ppt_path}")
    print(f"Certificate Width: {certificate.width}, Height: {certificate.height}")

    print("Certificates generated successfully!")
    
    # Open the output directory after generating certificates
    os.startfile(output_dir)

    '''time.sleep(3)
    sys.exit()'''
    